from flask import Flask, request, jsonify, send_from_directory
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
import re
import json
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain.schema import Document
from dotenv import load_dotenv

app = Flask(__name__, static_folder='../frontend', static_url_path='')

# Load environment variables
load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

def preprocess_text(text):
    """Preprocess the text by cleaning and formatting."""
    text = re.sub(r'\*\*(.*?)\*\*', '**\\1**', text)  # Ensure text contains formatting markers
    text = re.sub(r'_(.*?)_', '_\\1_', text)
    text = re.sub(r'\s+', ' ', text)  # Remove extra whitespace
    text = text.strip()  # Remove leading and trailing whitespace
    return text

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            page_text = page.extract_text() or ""
            text += page_text
    return text

def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    return text

def get_docx_text(docx_docs):
    text = ""
    for docx_file in docx_docs:
        doc = DocxDocument(docx_file)
        for para in doc.paragraphs:
            if para.text.strip():  # Check if paragraph text is not empty
                text += para.text + "\n"
    return text

def get_text_from_files(files):
    text = ""
    pdf_files = [f for f in files if f.filename.endswith('.pdf')]
    ppt_files = [f for f in files if f.filename.endswith('.pptx')]
    docx_files = [f for f in files if f.filename.endswith('.docx')]

    if pdf_files:
        text += get_pdf_text(pdf_files)
    if ppt_files:
        text += get_ppt_text(ppt_files)
    if docx_files:
        text += get_docx_text(docx_files)
        
    return text
    
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=5000, chunk_overlap=500)  # Adjust as needed
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    if not text_chunks:
        raise ValueError("Text chunks are empty")
    
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    # Debugging information
    print(f"Number of text chunks: {len(text_chunks)}")
    
    try:
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("faiss_index")
    except Exception as e:
        print(f"Error creating FAISS index: {e}")
        raise

def get_conversational_chain():
    prompt_template = """
    Answer the question in detail using the context provided. If the answer is not in the context, say "Answer is not available in the context."
    
    Context:
    {context}
    
    Question:
    {question}
    
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain


def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    
    # Debugging information
    print(f"Retrieved documents: {docs}")

    chain = get_conversational_chain()
    response = chain.invoke({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def summarize_pdf(text):
    prompt_template = """
    Summarize the following text in detail:\n\n
    Text:\n {context}\n

    Summary:
    """
    global model
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    document = Document(page_content=text)
    response = chain.invoke({"input_documents": [document]}, return_only_outputs=True)
    return response["output_text"]

@app.route('/')
def serve_index():
    return send_from_directory(app.static_folder, 'index.html')

@app.route('/ask_question', methods=['POST'])
def ask_question():
    try:
        data = json.loads(request.data)
        question = data.get('question', '')
        context = "Your document text or relevant context"
        context = preprocess_text(context)
        # Use the correct method to handle the question
        chain = get_conversational_chain()
        answer = chain.invoke({"input_documents": [context], "question": question}, return_only_outputs=True)
        return jsonify({'answer': preprocess_text(answer["output_text"])})
    except Exception as e:
        return jsonify({'error': str(e)})


@app.route('/submit_files', methods=['POST'])
def submit_files():
    files = request.files.getlist('files')
    raw_text = get_text_from_files(files)
    if raw_text:
        text_chunks = get_text_chunks(raw_text)
        if text_chunks:
            get_vector_store(text_chunks)
            return jsonify({'message': 'Files processed successfully'}), 200
        return jsonify({'error': 'No text chunks were generated.'}), 400
    return jsonify({'error': 'No text was extracted from the files.'}), 400

@app.route('/summarize', methods=['POST'])
def summarize():
    try:
        data = json.loads(request.data)
        files_data = data.get('files_data', '')
        # Ensure model is initialized
        if 'model' not in globals():
            get_conversational_chain()  # Ensure the model is initialized
        summary = summarize_pdf(files_data)
        return jsonify({'summary': preprocess_text(summary)})
    except Exception as e:
        return jsonify({'error': str(e)})

if __name__ == "__main__":
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)))
